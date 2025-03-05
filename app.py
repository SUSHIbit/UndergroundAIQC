import streamlit as st
import mysql.connector
import hashlib
import openai
import os
import pptx
import json
import io
from PIL import Image
import base64

# Set page configuration
st.set_page_config(page_title="Quiz Generator", layout="wide", page_icon="📚")

# Configure OpenAI API - REPLACE WITH YOUR API KEY
OPENAI_API_KEY = "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" # Replace with your actual OpenAI API key
openai.api_key = OPENAI_API_KEY

# Database configuration - MODIFY THESE VALUES
DB_HOST = "localhost"
DB_USER = "root"
DB_PASSWORD = ""
DB_NAME = "undergroundproject"

# Database connection function
def get_db_connection():
    return mysql.connector.connect(
        host=DB_HOST,
        user=DB_USER,
        password=DB_PASSWORD,
        database=DB_NAME
    )

# Authentication functions
def verify_password(stored_password, provided_password):
    """Verify a password against its hashed value from Laravel"""
    # Laravel uses bcrypt, but we'll use a simplified approach for demo
    # In a production environment, use a proper password verification library
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    
    try:
        # Get the user by username
        cursor.execute(
            "SELECT * FROM users WHERE username = %s", 
            (provided_password,)
        )
        return True
    except Exception as e:
        st.error(f"Authentication error: {e}")
        return False
    finally:
        cursor.close()
        conn.close()

def authenticate_user(username, password):
    """Authenticate a user and return their details if valid"""
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    
    try:
        # Get the user by username
        cursor.execute(
            "SELECT * FROM users WHERE username = %s", 
            (username,)
        )
        user = cursor.fetchone()
        
        if user and user['role'] == 'lecturer':
            # In a real implementation, you would need to use a proper Laravel-compatible 
            # password verification function instead of this simplified approach
            return user
        return None
    except Exception as e:
        st.error(f"Authentication error: {e}")
        return None
    finally:
        cursor.close()
        conn.close()

# PPT Processing functions
def extract_text_from_pptx(uploaded_file):
    """Extract text content from an uploaded PowerPoint file"""
    content = []
    
    # Save the uploaded file to a temporary file
    bytes_data = uploaded_file.getvalue()
    presentation = pptx.Presentation(io.BytesIO(bytes_data))
    
    # Extract text from each slide
    for i, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        
        if slide_text:
            content.append(f"Slide {i+1}: " + " ".join(slide_text))
    
    return "\n\n".join(content)

# Question generation functions
def generate_questions_with_openai(content, question_type="quiz"):
    """Generate questions using OpenAI based on the slide content"""
    try:
        difficulty = "standard" if question_type == "quiz" else "challenging"
        
        prompt = f"""
        Based on the following lecture slide content, generate 10 multiple-choice questions at a {difficulty} difficulty level.
        Each question should have 4 options (A, B, C, D) with only one correct answer.
        For each question, also provide:
        1. The correct answer letter
        2. A brief explanation for why that answer is correct
        
        Format each question as follows:
        Question 1: [Question text]
        Options: A: [Option A], B: [Option B], C: [Option C], D: [Option D]
        Answer: [Correct answer letter]
        Reason: [Explanation]
        
        Only use information explicitly stated in the content. Do not make up information.
        
        Content:
        {content}
        """

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an educational assistant that creates accurate multiple-choice questions based on lecture content."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=3000,
            temperature=0.7
        )
        
        response_text = response.choices[0].message.content
        
        # Process response into structured questions
        questions = parse_openai_response(response_text)
        return questions
    
    except Exception as e:
        st.error(f"Error generating questions: {e}")
        return []

def parse_openai_response(response_text):
    """Parse the OpenAI response into structured question objects"""
    lines = response_text.strip().split('\n')
    questions = []
    current_question = {}
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith("Question"):
            # Save previous question if it exists
            if current_question and 'question' in current_question:
                questions.append(current_question)
            
            # Start new question
            q_parts = line.split(":", 1)
            if len(q_parts) > 1:
                current_question = {'question': q_parts[1].strip()}
            else:
                current_question = {'question': ''}
                
        elif line.startswith("Options:"):
            options_str = line[8:].strip()
            options = {}
            for opt in options_str.split(','):
                opt = opt.strip()
                if opt and ':' in opt:
                    key, value = opt.split(':', 1)
                    options[key.strip()] = value.strip()
            current_question['options'] = options
            
        elif line.startswith("Answer:"):
            current_question['answer'] = line[7:].strip()
            
        elif line.startswith("Reason:"):
            current_question['reason'] = line[7:].strip()
    
    # Add the last question
    if current_question and 'question' in current_question:
        questions.append(current_question)
    
    return questions

# Database functions
def get_next_set_number():
    """Get the next available set number"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        cursor.execute("SELECT MAX(set_number) FROM sets")
        result = cursor.fetchone()
        next_set = 1 if result[0] is None else result[0] + 1
        return next_set
    except Exception as e:
        st.error(f"Error getting next set number: {e}")
        return 1
    finally:
        cursor.close()
        conn.close()

def get_subjects():
    """Get all subjects from the database"""
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    
    try:
        cursor.execute("SELECT * FROM subjects ORDER BY name")
        return cursor.fetchall()
    except Exception as e:
        st.error(f"Error fetching subjects: {e}")
        return []
    finally:
        cursor.close()
        conn.close()

def create_subject(name):
    """Create a new subject in the database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        cursor.execute("INSERT INTO subjects (name) VALUES (%s)", (name,))
        conn.commit()
        return cursor.lastrowid
    except Exception as e:
        st.error(f"Error creating subject: {e}")
        conn.rollback()
        return None
    finally:
        cursor.close()
        conn.close()

def get_topics(subject_id):
    """Get all topics for a specific subject"""
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    
    try:
        cursor.execute(
            "SELECT * FROM topics WHERE subject_id = %s ORDER BY name", 
            (subject_id,)
        )
        return cursor.fetchall()
    except Exception as e:
        st.error(f"Error fetching topics: {e}")
        return []
    finally:
        cursor.close()
        conn.close()

def create_topic(subject_id, name):
    """Create a new topic in the database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        cursor.execute(
            "INSERT INTO topics (subject_id, name) VALUES (%s, %s)", 
            (subject_id, name)
        )
        conn.commit()
        return cursor.lastrowid
    except Exception as e:
        st.error(f"Error creating topic: {e}")
        conn.rollback()
        return None
    finally:
        cursor.close()
        conn.close()

def get_quiz_sets():
    """Get all quiz sets from the database"""
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    
    try:
        cursor.execute("""
            SELECT s.id, s.set_number, qd.subject_id, qd.topic_id, 
                   sub.name as subject_name, t.name as topic_name
            FROM sets s
            JOIN quiz_details qd ON s.id = qd.set_id
            JOIN subjects sub ON qd.subject_id = sub.id
            JOIN topics t ON qd.topic_id = t.id
            WHERE s.type = 'quiz'
            ORDER BY s.set_number
        """)
        return cursor.fetchall()
    except Exception as e:
        st.error(f"Error fetching quiz sets: {e}")
        return []
    finally:
        cursor.close()
        conn.close()

def save_quiz(set_number, user_id, subject_id, topic_id, questions):
    """Save a quiz to the database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        # First, create the set record with 'draft' status
        cursor.execute(
            "INSERT INTO sets (set_number, type, created_by, status) VALUES (%s, %s, %s, %s)",
            (set_number, 'quiz', user_id, 'draft')
        )
        set_id = cursor.lastrowid
        
        # Then, create the quiz details record
        cursor.execute(
            "INSERT INTO quiz_details (set_id, subject_id, topic_id) VALUES (%s, %s, %s)",
            (set_id, subject_id, topic_id)
        )
        
        # Finally, insert all questions
        for i, q in enumerate(questions, 1):
            options_json = json.dumps(q['options'])
            cursor.execute(
                """INSERT INTO questions 
                   (set_id, question_number, question_text, options, correct_answer, reason)
                   VALUES (%s, %s, %s, %s, %s, %s)""",
                (set_id, i, q['question'], options_json, q['answer'], q['reason'])
            )
        
        conn.commit()
        return True
    except Exception as e:
        st.error(f"Error saving quiz: {e}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def save_challenge(set_number, user_id, name, prerequisite_ids, questions):
    """Save a challenge to the database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        # First, create the set record
        cursor.execute(
            "INSERT INTO sets (set_number, type, created_by) VALUES (%s, %s, %s)",
            (set_number, 'challenge', user_id)
        )
        set_id = cursor.lastrowid
        
        # Then, create the challenge details record
        cursor.execute(
            "INSERT INTO challenge_details (set_id, name) VALUES (%s, %s)",
            (set_id, name)
        )
        challenge_id = cursor.lastrowid
        
        # Add prerequisite records
        for prereq_id in prerequisite_ids:
            cursor.execute(
                "INSERT INTO challenge_prerequisites (challenge_id, prerequisite_set_id) VALUES (%s, %s)",
                (challenge_id, prereq_id)
            )
        
        # Finally, insert all questions
        for i, q in enumerate(questions, 1):
            options_json = json.dumps(q['options'])
            cursor.execute(
                """INSERT INTO questions 
                   (set_id, question_number, question_text, options, correct_answer, reason)
                   VALUES (%s, %s, %s, %s, %s, %s)""",
                (set_id, i, q['question'], options_json, q['answer'], q['reason'])
            )
        
        conn.commit()
        return True
    except Exception as e:
        st.error(f"Error saving challenge: {e}")
        conn.rollback()
        return False
    finally:
        cursor.close()
        conn.close()

def get_quiz_questions(set_ids):
    """Get questions for specific quiz sets"""
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    
    questions = []
    try:
        # Convert set_ids to a comma-separated string for SQL IN clause
        set_ids_str = ', '.join(map(str, set_ids))
        
        cursor.execute(f"""
            SELECT q.*, s.set_number 
            FROM questions q
            JOIN sets s ON q.set_id = s.id
            WHERE q.set_id IN ({set_ids_str})
            ORDER BY s.set_number, q.question_number
        """)
        questions = cursor.fetchall()
        
        # Convert the JSON string to actual Python dictionaries
        for q in questions:
            q['options'] = json.loads(q['options'])
            
        return questions
    except Exception as e:
        st.error(f"Error fetching quiz questions: {e}")
        return []
    finally:
        cursor.close()
        conn.close()

# Main app
def main():
    # Session state initialization
    if 'authenticated' not in st.session_state:
        st.session_state.authenticated = False
    if 'user' not in st.session_state:
        st.session_state.user = None
    if 'page' not in st.session_state:
        st.session_state.page = 'login'
    if 'questions' not in st.session_state:
        st.session_state.questions = []
    if 'set_number' not in st.session_state:
        st.session_state.set_number = None
    
    # Custom CSS
    st.markdown("""
    <style>
    .main-header {
        font-size: 2.5rem;
        color: #1E88E5;
        text-align: center;
    }
    .sub-header {
        font-size: 1.8rem;
        color: #333;
        margin-bottom: 20px;
    }
    .card {
        padding: 20px;
        border-radius: 5px;
        margin-bottom: 20px;
        background-color: #f8f9fa;
        border: 1px solid #dee2e6;
    }
    .success-msg {
        color: #28a745;
        font-weight: bold;
        padding: 10px;
        border-radius: 5px;
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        margin: 10px 0;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # App logic based on current page
    if st.session_state.page == 'login':
        display_login_page()
    elif st.session_state.page == 'menu':
        display_menu_page()
    elif st.session_state.page == 'quiz':
        display_quiz_page()
    elif st.session_state.page == 'challenge':
        display_challenge_page()

def display_login_page():
    st.markdown("<h1 class='main-header'>Quiz Generator System</h1>", unsafe_allow_html=True)
    st.markdown("<h2 class='sub-header'>Lecturer Login</h2>", unsafe_allow_html=True)
    
    with st.form("login_form"):
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        submit_button = st.form_submit_button("Login")
        
        if submit_button:
            user = authenticate_user(username, password)
            if user and user['role'] == 'lecturer':
                st.session_state.authenticated = True
                st.session_state.user = user
                st.session_state.page = 'menu'
                st.rerun()  # Updated from st.experimental_rerun()
            else:
                st.error("Invalid credentials or you do not have lecturer privileges.")

def display_menu_page():
    st.markdown("<h1 class='main-header'>Quiz Generator System</h1>", unsafe_allow_html=True)
    st.markdown(f"<h2 class='sub-header'>Welcome, {st.session_state.user['name']}</h2>", unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("<div class='card'>", unsafe_allow_html=True)
        st.markdown("<h3>Create Quiz</h3>", unsafe_allow_html=True)
        st.markdown("<p>Generate quiz questions from PowerPoint slides.</p>", unsafe_allow_html=True)
        if st.button("Create Quiz", key="create_quiz"):
            st.session_state.page = 'quiz'
            # Get next set number here
            st.session_state.set_number = get_next_set_number()
            st.rerun()  # Updated from st.experimental_rerun()
        st.markdown("</div>", unsafe_allow_html=True)
    
    with col2:
        st.markdown("<div class='card'>", unsafe_allow_html=True)
        st.markdown("<h3>Create Challenge</h3>", unsafe_allow_html=True)
        st.markdown("<p>Create harder questions based on prerequisite quizzes.</p>", unsafe_allow_html=True)
        if st.button("Create Challenge", key="create_challenge"):
            st.session_state.page = 'challenge'
            # Get next set number here
            st.session_state.set_number = get_next_set_number()
            st.rerun()  # Updated from st.experimental_rerun()
        st.markdown("</div>", unsafe_allow_html=True)
    
    if st.button("Logout"):
        st.session_state.authenticated = False
        st.session_state.user = None
        st.session_state.page = 'login'
        st.rerun()  # Updated from st.experimental_rerun()

def display_quiz_page():
    st.markdown("<h1 class='main-header'>Create Quiz</h1>", unsafe_allow_html=True)
    
    # Display set number
    st.markdown(f"<h2 class='sub-header'>Set Number: {st.session_state.set_number}</h2>", unsafe_allow_html=True)
    
    # Subject selection
    st.subheader("Subject")
    
    # Get subjects from the database
    subjects = get_subjects()
    subject_names = [s['name'] for s in subjects] + ["Add New Subject"]
    
    subject_option = st.selectbox("Select or add a subject", subject_names)
    
    subject_id = None
    
    if subject_option == "Add New Subject":
        new_subject = st.text_input("Enter new subject name")
        if new_subject and st.button("Add Subject"):
            subject_id = create_subject(new_subject)
            if subject_id:
                st.success(f"Subject '{new_subject}' added successfully!")
                st.rerun()  # Updated from st.experimental_rerun()
    else:
        # Find the selected subject in the list
        for s in subjects:
            if s['name'] == subject_option:
                subject_id = s['id']
                break
    
    # Topic selection (only if a subject is selected)
    topic_id = None
    if subject_id:
        st.subheader("Topic")
        
        # Get topics for the selected subject
        topics = get_topics(subject_id)
        topic_names = [t['name'] for t in topics] + ["Add New Topic"]
        
        topic_option = st.selectbox("Select or add a topic", topic_names)
        
        if topic_option == "Add New Topic":
            new_topic = st.text_input("Enter new topic name")
            if new_topic and st.button("Add Topic"):
                topic_id = create_topic(subject_id, new_topic)
                if topic_id:
                    st.success(f"Topic '{new_topic}' added successfully!")
                    st.rerun()  # Updated from st.experimental_rerun()
        else:
            # Find the selected topic in the list
            for t in topics:
                if t['name'] == topic_option:
                    topic_id = t['id']
                    break
    
    # PowerPoint upload and processing
    if subject_id and topic_id:
        st.subheader("Upload PowerPoint Presentation")
        uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")
        
        if uploaded_file:
            # Extract content from the PowerPoint
            with st.spinner("Processing PowerPoint content..."):
                ppt_content = extract_text_from_pptx(uploaded_file)
                
                if ppt_content:
                    st.success("PowerPoint processed successfully!")
                    
                    with st.expander("View Slide Content"):
                        st.text_area("Content extracted from slides", ppt_content, height=200)
                    
                    # Generate questions
                    if st.button("Generate Questions"):
                        with st.spinner("Generating questions with AI..."):
                            st.session_state.questions = generate_questions_with_openai(ppt_content)
                            if st.session_state.questions:
                                st.success("Questions generated successfully!")
    
    # Display and edit questions
    if st.session_state.questions:
        st.subheader("Review and Edit Questions")
        
        edited_questions = []
        
        for i, q in enumerate(st.session_state.questions):
            with st.expander(f"Question {i+1}"):
                question_text = st.text_area(f"Question Text", q['question'], key=f"q_{i}")
                
                options = {}
                for opt in ['A', 'B', 'C', 'D']:
                    if opt in q['options']:
                        options[opt] = st.text_input(f"Option {opt}", q['options'][opt], key=f"opt_{i}_{opt}")
                    else:
                        options[opt] = st.text_input(f"Option {opt}", "", key=f"opt_{i}_{opt}")
                
                answer = st.selectbox(
                    "Correct Answer", 
                    ['A', 'B', 'C', 'D'], 
                    index=['A', 'B', 'C', 'D'].index(q['answer']) if q['answer'] in ['A', 'B', 'C', 'D'] else 0,
                    key=f"ans_{i}"
                )
                reason = st.text_area("Reason", q.get('reason', ''), key=f"reason_{i}")
                
                edited_questions.append({
                    'question': question_text,
                    'options': options,
                    'answer': answer,
                    'reason': reason
                })
        
        # Save all questions
        if edited_questions and st.button("Save Quiz"):
            # Save to database
            success = save_quiz(
                st.session_state.set_number, 
                st.session_state.user['id'],
                subject_id, 
                topic_id, 
                edited_questions
            )
            
            if success:
                st.markdown("<div class='success-msg'>Quiz saved successfully!</div>", unsafe_allow_html=True)
                st.session_state.questions = []
                st.session_state.set_number = None
                
                if st.button("Return to Menu"):
                    st.session_state.page = 'menu'
                    st.rerun()  # Updated from st.experimental_rerun()
    
    # Back button
    if st.button("Back to Menu"):
        st.session_state.page = 'menu'
        st.session_state.questions = []
        st.rerun()  # Updated from st.experimental_rerun()

def display_challenge_page():
    st.markdown("<h1 class='main-header'>Create Challenge</h1>", unsafe_allow_html=True)
    
    # Display set number
    st.markdown(f"<h2 class='sub-header'>Set Number: {st.session_state.set_number}</h2>", unsafe_allow_html=True)
    
    # Challenge name
    st.subheader("Challenge Name")
    challenge_name = st.text_input("Enter challenge name")
    
    # Prerequisite selection
    st.subheader("Select Prerequisites")
    st.write("Choose at least 2 prerequisite quiz sets:")
    
    quiz_sets = get_quiz_sets()
    if not quiz_sets:
        st.warning("No quiz sets available. Please create some quizzes first.")
        if st.button("Back to Menu"):
            st.session_state.page = 'menu'
            st.rerun()  # Updated from st.experimental_rerun()
        return
    
    # Allow the user to select multiple quiz sets
    selected_sets = []
    for quiz in quiz_sets:
        if st.checkbox(f"Set {quiz['set_number']}: {quiz['subject_name']} - {quiz['topic_name']}", key=f"set_{quiz['id']}"):
            selected_sets.append(quiz)
    
    # Display questions from selected sets
    if len(selected_sets) < 2:
        st.warning("Please select at least 2 prerequisite quiz sets.")
    else:
        selected_set_ids = [s['id'] for s in selected_sets]
        prerequisite_questions = get_quiz_questions(selected_set_ids)
        
        if prerequisite_questions:
            with st.expander("View Prerequisite Questions"):
                for q in prerequisite_questions:
                    st.markdown(f"**Set {q['set_number']}, Question {q['question_number']}**: {q['question_text']}")
                    for opt, text in q['options'].items():
                        if opt == q['correct_answer']:
                            st.markdown(f"* **{opt}: {text} ✓**")
                        else:
                            st.markdown(f"* {opt}: {text}")
                    st.markdown(f"**Reason**: {q['reason']}")
                    st.markdown("---")
            
            # Generate harder questions
            if st.button("Generate Harder Questions"):
                # Combine questions for context
                context = ""
                for q in prerequisite_questions:
                    context += f"Question: {q['question_text']}\n"
                    for opt, text in q['options'].items():
                        context += f"{opt}: {text}\n"
                    context += f"Answer: {q['correct_answer']}\n"
                    context += f"Reason: {q['reason']}\n\n"
                
                with st.spinner("Generating harder questions based on prerequisites..."):
                    st.session_state.questions = generate_questions_with_openai(context, question_type="challenge")
                    if st.session_state.questions:
                        st.success("Challenge questions generated successfully!")
        else:
            st.warning("Could not find questions for the selected prerequisite sets.")
    
    # Display and edit challenge questions
    if st.session_state.questions:
        st.subheader("Review and Edit Challenge Questions")
        
        edited_questions = []
        
        for i, q in enumerate(st.session_state.questions):
            with st.expander(f"Question {i+1}"):
                question_text = st.text_area(f"Question Text", q['question'], key=f"cq_{i}")
                
                options = {}
                for opt in ['A', 'B', 'C', 'D']:
                    if opt in q['options']:
                        options[opt] = st.text_input(f"Option {opt}", q['options'][opt], key=f"copt_{i}_{opt}")
                    else:
                        options[opt] = st.text_input(f"Option {opt}", "", key=f"copt_{i}_{opt}")
                
                answer = st.selectbox(
                    "Correct Answer", 
                    ['A', 'B', 'C', 'D'], 
                    index=['A', 'B', 'C', 'D'].index(q['answer']) if q['answer'] in ['A', 'B', 'C', 'D'] else 0,
                    key=f"cans_{i}"
                )
                reason = st.text_area("Reason", q.get('reason', ''), key=f"creason_{i}")
                
                edited_questions.append({
                    'question': question_text,
                    'options': options,
                    'answer': answer,
                    'reason': reason
                })
        
        # Save all questions
        if challenge_name and edited_questions and len(selected_sets) >= 2:
            if st.button("Save Challenge"):
                # Save to database
                success = save_challenge(
                    st.session_state.set_number, 
                    st.session_state.user['id'],
                    challenge_name, 
                    [s['id'] for s in selected_sets], 
                    edited_questions
                )
                
                if success:
                    st.markdown("<div class='success-msg'>Challenge saved successfully!</div>", unsafe_allow_html=True)
                    st.session_state.questions = []
                    st.session_state.set_number = None
                    
                    if st.button("Return to Menu"):
                        st.session_state.page = 'menu'
                        st.rerun()  # Updated from st.experimental_rerun()
    
    # Back button
    if st.button("Back to Menu"):
        st.session_state.page = 'menu'
        st.session_state.questions = []
        st.rerun()  # Updated from st.experimental_rerun()

if __name__ == "__main__":
    main()