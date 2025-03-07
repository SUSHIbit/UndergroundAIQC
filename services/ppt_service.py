import io
import pptx
import streamlit as st

def extract_text_from_pptx(uploaded_file):
    """
    Extract text content from an uploaded PowerPoint file
    
    Args:
        uploaded_file: The uploaded file object from Streamlit
    
    Returns:
        str: The extracted text content
    """
    content = []
    
    try:
        # Read the uploaded file
        bytes_data = uploaded_file.getvalue()
        presentation = pptx.Presentation(io.BytesIO(bytes_data))
        
        # Extract text from each slide
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            if slide_text:
                content.append(f"Slide {i+1}: " + " ".join(slide_text))
        
        return "\n\n".join(content)
    except Exception as e:
        st.error(f"Error processing PowerPoint file: {e}")
        return ""